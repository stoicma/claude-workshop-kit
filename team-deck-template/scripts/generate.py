"""
Render a McKinsey-style deck plan (markdown) to a .pptx with real consulting visuals.

Usage:
    python3 generate.py <input.md> <output.pptx> [--template path/to/template.pptx]

Visual specs supported (parsed from the `**Visual:**` line):
- bar:           bar chart (single series)
- line:          line chart (1 series)
- column-grouped: grouped column chart (multi-series)
- pie:           pie chart
- matrix-2x3:    2-row by 3-column matrix
- 3-card:        three-card grid
- comparison:    two-column side-by-side
- roadmap:       horizontal phased timeline
- risk-table:    risk and mitigation table
- decision-table: decision, owner, deadline table
- none:          no visual (used by title and divider layouts)

See SKILL.md and rules/visuals.md for the spec format.
"""

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
except ImportError:
    sys.stderr.write("python-pptx is not installed. Run: pip install python-pptx\n")
    sys.exit(1)


# ---------- Color palette ----------
NAVY = RGBColor(0x0A, 0x25, 0x40)
ACCENT = RGBColor(0xE7, 0x6F, 0x00)
ACCENT_LIGHT = RGBColor(0xFB, 0xC4, 0x8F)
GREY = RGBColor(0x66, 0x66, 0x66)
GREY_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
GREY_BG = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------- Markdown parser ----------
def parse_markdown(md_text):
    slides = []
    deck_meta = {"title": "", "audience": "", "length": "", "scr": []}
    current = None
    in_meta = True

    title_re = re.compile(r"^# (.+)$")
    slide_re = re.compile(r"^## Slide (\d+): (.+)$")
    field_re = re.compile(r"^\*\*([^*]+):\*\*\s*(.*)$")

    for line in md_text.splitlines():
        m = title_re.match(line)
        if m and in_meta:
            deck_meta["title"] = m.group(1).strip()
            continue
        m = slide_re.match(line)
        if m:
            in_meta = False
            if current is not None:
                slides.append(current)
            current = {
                "number": int(m.group(1)),
                "title": m.group(2).strip(),
                "layout": "content",
                "body": [],
                "visual": "none",
                "source": "",
            }
            continue
        if in_meta:
            fm = field_re.match(line)
            if fm:
                key = fm.group(1).strip().lower()
                val = fm.group(2).strip()
                if key == "audience":
                    deck_meta["audience"] = val
                elif key in ("length target", "length"):
                    deck_meta["length"] = val
                elif key == "storyline (scr)":
                    pass
            elif line.startswith("- **"):
                deck_meta["scr"].append(line.lstrip("- "))
            continue
        if current is None:
            continue
        fm = field_re.match(line)
        if fm:
            field = fm.group(1).strip().lower()
            value = fm.group(2).strip()
            if field == "layout":
                current["layout"] = value or "content"
            elif field == "body":
                if value:
                    current["body"].append(value)
            elif field == "visual":
                current["visual"] = value
            elif field in ("source footer", "source"):
                current["source"] = value
            continue
        bullet = line.lstrip()
        if bullet.startswith("- "):
            current["body"].append(bullet[2:])

    if current is not None:
        slides.append(current)
    return deck_meta, slides


# ---------- Text helpers ----------
def set_text(shape, text, font_size=14, bold=False, color=NAVY, italic=False, font_name="Helvetica", align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def add_text(slide, text, x, y, w, h, **kwargs):
    box = slide.shapes.add_textbox(x, y, w, h)
    set_text(box, text, **kwargs)
    return box


def add_filled_box(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_line(slide, x, y, w):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


# ---------- Slide-level layouts ----------
def add_title_slide(prs, slide_data, deck_meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_filled_box(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    add_filled_box(slide, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), NAVY)
    add_filled_box(slide, Inches(0.5), Inches(2.4), Inches(0.15), Inches(0.6), ACCENT)

    add_text(slide, slide_data["title"],
             Inches(0.85), Inches(2.3), Inches(11.8), Inches(1.6),
             font_size=36, bold=True, color=NAVY)

    sub_lines = []
    if deck_meta.get("audience"):
        sub_lines.append(f"Prepared for: {deck_meta['audience']}")
    if slide_data["body"]:
        sub_lines.extend(slide_data["body"])
    if sub_lines:
        add_text(slide, "  |  ".join(sub_lines),
                 Inches(0.85), Inches(4.2), Inches(11.8), Inches(0.8),
                 font_size=14, color=GREY)


def add_divider_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_box(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)

    add_filled_box(slide, Inches(0.5), Inches(3.0), Inches(0.6), Inches(0.12), ACCENT)
    add_text(slide, slide_data["title"],
             Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.5),
             font_size=32, bold=True, color=WHITE)

    if slide_data["body"]:
        add_text(slide, " ".join(slide_data["body"]),
                 Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.5),
                 font_size=18, color=ACCENT_LIGHT, italic=True)


def add_content_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_filled_box(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)

    # Header
    add_text(slide, slide_data["title"],
             Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.9),
             font_size=22, bold=True, color=NAVY)
    add_accent_line(slide, Inches(0.5), Inches(1.35), Inches(0.8))

    # Body region
    body_top = Inches(1.6)
    body_height = Inches(5.2)
    visual_spec = slide_data["visual"].strip()

    if visual_spec.lower() == "none" or not visual_spec:
        render_text_only(slide, slide_data, body_top, body_height)
    else:
        render_visual(slide, visual_spec, slide_data, body_top, body_height)

    # Footer source
    if slide_data["source"]:
        add_text(slide, slide_data["source"],
                 Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
                 font_size=8, color=GREY, italic=True)


def render_text_only(slide, slide_data, top, height):
    if not slide_data["body"]:
        return
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(slide_data["body"]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "•   " + item
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.name = "Helvetica"
            run.font.size = Pt(16)
            run.font.color.rgb = NAVY


# ---------- Visual dispatcher ----------
def render_visual(slide, spec, slide_data, top, height):
    spec_type, _, spec_body = spec.partition(":")
    spec_type = spec_type.strip().lower()
    spec_body = spec_body.strip()

    body_text = " ".join(slide_data["body"]) if slide_data["body"] else ""
    body_box_top = top
    body_box_height = Inches(1.0)
    visual_top = top + Inches(1.1)
    visual_height = height - Inches(1.1)

    if slide_data["body"]:
        body_box = slide.shapes.add_textbox(Inches(0.5), body_box_top, Inches(12.3), body_box_height)
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(slide_data["body"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "•  " + item
            p.space_after = Pt(2)
            for run in p.runs:
                run.font.name = "Helvetica"
                run.font.size = Pt(12)
                run.font.color.rgb = NAVY
    else:
        visual_top = top
        visual_height = height

    handlers = {
        "bar": render_bar,
        "line": render_line,
        "column-grouped": render_column_grouped,
        "pie": render_pie,
        "matrix-2x3": render_matrix,
        "3-card": render_three_card,
        "comparison": render_comparison,
        "roadmap": render_roadmap,
        "risk-table": render_risk_table,
        "decision-table": render_decision_table,
    }
    fn = handlers.get(spec_type)
    if fn:
        fn(slide, spec_body, visual_top, visual_height)
    else:
        sys.stderr.write(f"Unknown visual type: {spec_type}\n")


# ---------- Visual renderers ----------
def render_bar(slide, body, top, height):
    title, _, data = body.partition("|")
    chart_data = CategoryChartData()
    pairs = [p.strip() for p in data.split(",") if "=" in p]
    cats = [p.split("=", 1)[0].strip() for p in pairs]
    vals = [float(p.split("=", 1)[1].strip()) for p in pairs]
    chart_data.categories = cats
    chart_data.add_series(title.strip() or "Series 1", vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.5), top, Inches(10.3), height,
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.has_title = bool(title.strip())
    if title.strip():
        chart.chart_title.text_frame.text = title.strip()
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = NAVY
            run.font.name = "Helvetica"

    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.color.rgb = NAVY

    series = chart.series[0]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def render_line(slide, body, top, height):
    parts = [p.strip() for p in body.split("|")]
    title = parts[0] if len(parts) > 0 else ""
    data_str = parts[1] if len(parts) > 1 else ""
    series_name = "Series"
    for extra in parts[2:]:
        if extra.lower().startswith("series="):
            series_name = extra.split("=", 1)[1].strip()

    chart_data = CategoryChartData()
    pairs = [p.strip() for p in data_str.split(",") if "=" in p]
    cats = [p.split("=", 1)[0].strip() for p in pairs]
    vals = [float(p.split("=", 1)[1].strip()) for p in pairs]
    chart_data.categories = cats
    chart_data.add_series(series_name, vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(1.5), top, Inches(10.3), height,
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = NAVY
            run.font.name = "Helvetica"

    series = chart.series[0]
    line = series.format.line
    line.color.rgb = ACCENT
    line.width = Pt(2.5)


def render_column_grouped(slide, body, top, height):
    parts = [p.strip() for p in body.split("|")]
    title = parts[0] if len(parts) > 0 else ""
    series_specs = parts[1:]

    chart_data = CategoryChartData()
    cats_seen = []
    series_data = []

    for spec in series_specs:
        if ":" not in spec:
            continue
        sname, _, sdata = spec.partition(":")
        sname = sname.strip()
        if sname.lower().startswith("series"):
            sname = sname.split("=", 1)[1].strip() if "=" in sname else sname
        pairs = [p.strip() for p in sdata.split(",") if "=" in p]
        cats = [p.split("=", 1)[0].strip() for p in pairs]
        vals = [float(p.split("=", 1)[1].strip()) for p in pairs]
        if not cats_seen:
            cats_seen = cats
        series_data.append((sname, vals))

    chart_data.categories = cats_seen
    for sname, vals in series_data:
        chart_data.add_series(sname, vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.5), top, Inches(10.3), height,
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = NAVY

    colors = [NAVY, ACCENT, GREY]
    for i, ser in enumerate(chart.series):
        fill = ser.format.fill
        fill.solid()
        fill.fore_color.rgb = colors[i % 3]


def render_pie(slide, body, top, height):
    title, _, data = body.partition("|")
    chart_data = CategoryChartData()
    pairs = [p.strip() for p in data.split(",") if "=" in p]
    cats = [p.split("=", 1)[0].strip() for p in pairs]
    vals = [float(p.split("=", 1)[1].strip()) for p in pairs]
    chart_data.categories = cats
    chart_data.add_series(title.strip() or "Series", vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(2.5), top, Inches(8.3), height,
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.has_title = bool(title.strip())
    if title.strip():
        chart.chart_title.text_frame.text = title.strip()


def render_matrix(slide, body, top, height):
    rows = [r.strip() for r in body.split("|") if r.strip()]
    parsed = []
    for row in rows:
        if "=" not in row:
            continue
        label, _, cells = row.partition("=")
        cell_list = [c.strip() for c in cells.split(";")]
        parsed.append((label.strip(), cell_list))

    if not parsed:
        return

    n_rows = len(parsed)
    n_cols = max(len(c) for _, c in parsed)
    matrix_left = Inches(1.0)
    matrix_top = top
    label_col_w = Inches(2.0)
    cell_w = (Inches(11.3) - label_col_w) / n_cols
    cell_h = height / n_rows

    accent_cell = (n_rows - 1, n_cols - 1)

    for r, (label, cells) in enumerate(parsed):
        # row label
        lbox = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            matrix_left, matrix_top + cell_h * r, label_col_w, cell_h,
        )
        lbox.fill.solid()
        lbox.fill.fore_color.rgb = NAVY
        lbox.line.fill.background()
        lbox.shadow.inherit = False
        set_text(lbox, label,
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        lbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        for c in range(n_cols):
            text = cells[c] if c < len(cells) else ""
            is_accent = (r == accent_cell[0] and c == accent_cell[1])
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                matrix_left + label_col_w + cell_w * c,
                matrix_top + cell_h * r,
                cell_w, cell_h,
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT if is_accent else GREY_BG
            cell.line.color.rgb = WHITE
            cell.line.width = Pt(2)
            cell.shadow.inherit = False
            set_text(cell, text,
                     font_size=11,
                     color=WHITE if is_accent else NAVY,
                     align=PP_ALIGN.CENTER,
                     bold=is_accent)
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def render_three_card(slide, body, top, height):
    cards = [c.strip() for c in body.split("|") if c.strip()]
    cards = cards[:3]
    if not cards:
        return

    n = len(cards)
    margin = Inches(0.5)
    gap = Inches(0.3)
    total_w = Inches(13.333) - 2 * margin - (n - 1) * gap
    card_w = total_w / n
    card_h = min(height, Inches(4.5))
    card_y = top + (height - card_h) / 2

    for i, card in enumerate(cards):
        title, _, body_text = card.partition(":")
        x = margin + i * (card_w + gap)

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, card_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = GREY_BG
        bg.line.fill.background()
        bg.shadow.inherit = False

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        x, card_y, card_w, Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT
        accent.line.fill.background()
        accent.shadow.inherit = False

        num = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.3), card_y + Inches(0.4),
                                     Inches(0.6), Inches(0.6))
        num.fill.solid()
        num.fill.fore_color.rgb = NAVY
        num.line.fill.background()
        num.shadow.inherit = False
        set_text(num, str(i + 1), font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        num.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_box = slide.shapes.add_textbox(
            x + Inches(0.3), card_y + Inches(1.2),
            card_w - Inches(0.6), Inches(1.2),
        )
        set_text(title_box, title.strip(), font_size=15, bold=True, color=NAVY)

        if body_text.strip():
            body_box = slide.shapes.add_textbox(
                x + Inches(0.3), card_y + Inches(2.5),
                card_w - Inches(0.6), card_h - Inches(2.7),
            )
            set_text(body_box, body_text.strip(), font_size=12, color=GREY)


def render_comparison(slide, body, top, height):
    sides = [s.strip() for s in body.split("|") if s.strip()]
    sides = sides[:2]
    if len(sides) < 2:
        return

    margin = Inches(0.5)
    gap = Inches(0.4)
    col_w = (Inches(13.333) - 2 * margin - gap) / 2
    col_h = min(height, Inches(4.5))
    col_y = top + (height - col_h) / 2

    for i, side in enumerate(sides):
        title, _, bullets_text = side.partition(":")
        x = margin + i * (col_w + gap)

        is_left = (i == 0)
        bg_color = GREY_BG if is_left else ACCENT_LIGHT

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, col_y, col_w, col_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()
        bg.shadow.inherit = False

        title_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             x, col_y, col_w, Inches(0.7))
        title_strip.fill.solid()
        title_strip.fill.fore_color.rgb = NAVY if is_left else ACCENT
        title_strip.line.fill.background()
        title_strip.shadow.inherit = False
        set_text(title_strip, title.strip(),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        title_strip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        bullets = [b.strip() for b in bullets_text.split(";") if b.strip()]
        body_box = slide.shapes.add_textbox(
            x + Inches(0.2), col_y + Inches(0.9),
            col_w - Inches(0.4), col_h - Inches(1.0),
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(bullets):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "•  " + bullet
            p.space_after = Pt(6)
            for run in p.runs:
                run.font.name = "Helvetica"
                run.font.size = Pt(12)
                run.font.color.rgb = NAVY


def render_roadmap(slide, body, top, height):
    phases = [p.strip() for p in body.split("|") if p.strip()]
    n = len(phases)
    if n == 0:
        return

    margin = Inches(0.5)
    gap = Inches(0.2)
    arrow_w = Inches(0.4)
    available_w = Inches(13.333) - 2 * margin - (n - 1) * gap - (n - 1) * arrow_w
    box_w = available_w / n
    box_h = min(Inches(2.5), height - Inches(0.5))
    y = top + (height - box_h) / 2

    for i, phase in enumerate(phases):
        # phase format: "name=duration: outcome" or "name: outcome"
        name_part, _, outcome = phase.partition(":")
        if "=" in name_part:
            name, _, duration = name_part.partition("=")
        else:
            name, duration = name_part, ""
        x = margin + i * (box_w + arrow_w + gap)

        is_last = (i == n - 1)
        is_highlighted = (i == 1) if n >= 3 else (i == 0)
        bg_color = ACCENT if is_highlighted else NAVY

        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.fill.background()
        box.shadow.inherit = False

        name_box = slide.shapes.add_textbox(
            x + Inches(0.15), y + Inches(0.2),
            box_w - Inches(0.3), Inches(0.5),
        )
        set_text(name_box, name.strip(), font_size=14, bold=True, color=WHITE)

        if duration.strip():
            dur_box = slide.shapes.add_textbox(
                x + Inches(0.15), y + Inches(0.7),
                box_w - Inches(0.3), Inches(0.4),
            )
            set_text(dur_box, duration.strip(), font_size=10, color=ACCENT_LIGHT if is_highlighted else GREY_LIGHT, italic=True)

        if outcome.strip():
            out_box = slide.shapes.add_textbox(
                x + Inches(0.15), y + Inches(1.15),
                box_w - Inches(0.3), box_h - Inches(1.3),
            )
            set_text(out_box, outcome.strip(), font_size=11, color=WHITE)

        if not is_last:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + box_w + Inches(0.05), y + box_h / 2 - Inches(0.2),
                arrow_w, Inches(0.4),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GREY_LIGHT
            arrow.line.fill.background()
            arrow.shadow.inherit = False


def render_risk_table(slide, body, top, height):
    rows = [r.strip() for r in body.split("|") if r.strip()]
    pairs = []
    for row in rows:
        if "=" not in row:
            continue
        risk, _, mitig = row.partition("=")
        pairs.append((risk.strip(), mitig.strip()))
    if not pairs:
        return

    rows_n = len(pairs) + 1
    cols_n = 2
    margin = Inches(1.0)
    table_w = Inches(11.3)
    table_h = min(height, Inches(0.7) * rows_n + Inches(0.3))
    table_top = top + (height - table_h) / 2

    table_shape = slide.shapes.add_table(rows_n, cols_n, margin, table_top, table_w, table_h)
    table = table_shape.table
    table.columns[0].width = Inches(4.5)
    table.columns[1].width = Inches(6.8)

    headers = ["Risk", "Mitigation"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.text = h
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.name = "Helvetica"
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = WHITE

    for r, (risk, mitig) in enumerate(pairs, start=1):
        for c, val in enumerate([risk, mitig]):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else GREY_BG
            tf = cell.text_frame
            tf.text = val
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.name = "Helvetica"
                    run.font.size = Pt(11)
                    run.font.color.rgb = NAVY


def render_decision_table(slide, body, top, height):
    rows = [r.strip() for r in body.split("|") if r.strip()]
    parsed = []
    for row in rows:
        parts = [p.strip() for p in row.split(":")]
        while len(parts) < 3:
            parts.append("")
        parsed.append(parts[:3])
    if not parsed:
        return

    rows_n = len(parsed) + 1
    margin = Inches(1.0)
    table_w = Inches(11.3)
    table_h = min(height, Inches(0.6) * rows_n + Inches(0.3))
    table_top = top + (height - table_h) / 2

    table_shape = slide.shapes.add_table(rows_n, 3, margin, table_top, table_w, table_h)
    table = table_shape.table
    table.columns[0].width = Inches(6.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.3)

    headers = ["Decision", "Owner", "Deadline"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.text = h
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.name = "Helvetica"
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = WHITE

    for r, parts in enumerate(parsed, start=1):
        for c, val in enumerate(parts):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else GREY_BG
            tf = cell.text_frame
            tf.text = val
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.name = "Helvetica"
                    run.font.size = Pt(11)
                    run.font.color.rgb = NAVY


# ---------- Main ----------
def build_deck(deck_meta, slides, template_path=None):
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    for slide in slides:
        layout = slide["layout"].lower()
        if "title" in layout or slide["number"] == 1:
            add_title_slide(prs, slide, deck_meta)
        elif "divider" in layout or layout == "section":
            add_divider_slide(prs, slide)
        else:
            add_content_slide(prs, slide)
    return prs


def main():
    parser = argparse.ArgumentParser(description="Render a McKinsey-style deck plan to .pptx")
    parser.add_argument("input")
    parser.add_argument("output")
    parser.add_argument("--template", default=None)
    args = parser.parse_args()

    md_text = Path(args.input).read_text(encoding="utf-8")
    deck_meta, slides = parse_markdown(md_text)
    if not slides:
        sys.stderr.write("No slides parsed.\n")
        sys.exit(1)

    prs = build_deck(deck_meta, slides, template_path=args.template)
    prs.save(args.output)
    print(f"Wrote {len(slides)} slides to {args.output}")


if __name__ == "__main__":
    main()
